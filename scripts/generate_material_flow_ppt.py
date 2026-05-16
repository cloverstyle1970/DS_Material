"""
DS_Material 자재 흐름 PPT 생성 스크립트
실행: python scripts/generate_material_flow_ppt.py
산출: docs/자재흐름_DS_Material.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

# ──────────────────────────── 색상 팔레트 ────────────────────────────
NAVY      = RGBColor(0x1E, 0x3A, 0x8A)   # 본문 강조
BLUE      = RGBColor(0x25, 0x63, 0xEB)   # 일반 흐름 강조
AMBER     = RGBColor(0xD9, 0x77, 0x06)   # 견적 흐름 강조
EMERALD   = RGBColor(0x05, 0x96, 0x69)
ROSE      = RGBColor(0xE1, 0x1D, 0x48)
SLATE_900 = RGBColor(0x0F, 0x17, 0x2A)
SLATE_700 = RGBColor(0x33, 0x41, 0x55)
SLATE_500 = RGBColor(0x64, 0x74, 0x8B)
SLATE_200 = RGBColor(0xE2, 0xE8, 0xF0)
SLATE_100 = RGBColor(0xF1, 0xF5, 0xF9)
SLATE_50  = RGBColor(0xF8, 0xFA, 0xFC)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

# ──────────────────────────── 헬퍼 ────────────────────────────
WIDE_W, WIDE_H = Inches(13.333), Inches(7.5)   # 16:9

def add_text(slide, left, top, width, height, text, *,
             size=14, bold=False, color=SLATE_900, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, font="맑은 고딕"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top  = tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb

def add_rect(slide, left, top, width, height, *,
             fill=WHITE, line=SLATE_200, line_w=0.75, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    s = slide.shapes.add_shape(shape, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = line
    s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s

def step_box(slide, left, top, width, height, num, title, body, *, accent=BLUE):
    """단계 박스: [번호 뱃지] + 제목 + 설명"""
    add_rect(slide, left, top, width, height, fill=WHITE, line=SLATE_200, line_w=1.0)
    # 번호 뱃지
    badge_size = Inches(0.5)
    add_rect(slide, left + Inches(0.18), top + Inches(0.18),
             badge_size, badge_size, fill=accent, line=accent, shape=MSO_SHAPE.OVAL)
    add_text(slide, left + Inches(0.18), top + Inches(0.18),
             badge_size, badge_size, str(num),
             size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 제목
    add_text(slide, left + Inches(0.78), top + Inches(0.12),
             width - Inches(0.85), Inches(0.4), title,
             size=12, bold=True, color=SLATE_900)
    # 본문
    add_text(slide, left + Inches(0.78), top + Inches(0.55),
             width - Inches(0.85), height - Inches(0.65), body,
             size=10, color=SLATE_700)

def arrow_right(slide, left, top, width=Inches(0.4), height=Inches(0.3), color=SLATE_500):
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def arrow_down(slide, left, top, width=Inches(0.3), height=Inches(0.4), color=SLATE_500):
    s = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def header_bar(slide, title, subtitle, *, accent=BLUE):
    """상단 헤더 바 — 모든 본문 슬라이드 공통"""
    add_rect(slide, Inches(0), Inches(0), WIDE_W, Inches(0.9),
             fill=SLATE_50, line=SLATE_50, shape=MSO_SHAPE.RECTANGLE)
    # 좌측 색상 강조 띠
    add_rect(slide, Inches(0), Inches(0), Inches(0.18), Inches(0.9),
             fill=accent, line=accent, shape=MSO_SHAPE.RECTANGLE)
    add_text(slide, Inches(0.4), Inches(0.12), Inches(12), Inches(0.45),
             title, size=22, bold=True, color=SLATE_900)
    add_text(slide, Inches(0.4), Inches(0.52), Inches(12), Inches(0.35),
             subtitle, size=12, color=SLATE_500)

def footer(slide, n, total):
    add_text(slide, Inches(0.4), Inches(7.1), Inches(12), Inches(0.3),
             f"DS_Material — 자재 흐름  ·  {n} / {total}",
             size=9, color=SLATE_500)

# ──────────────────────────── 슬라이드 작성 ────────────────────────────

prs = Presentation()
prs.slide_width  = WIDE_W
prs.slide_height = WIDE_H
BLANK = prs.slide_layouts[6]   # 빈 레이아웃

TOTAL_SLIDES = 18

def new_slide():
    return prs.slides.add_slide(BLANK)

# ─── 1. 표지 ───
s = new_slide()
add_rect(s, Inches(0), Inches(0), WIDE_W, WIDE_H, fill=NAVY, line=NAVY,
         shape=MSO_SHAPE.RECTANGLE)
add_rect(s, Inches(0), Inches(5.2), WIDE_W, Inches(2.3), fill=SLATE_900, line=SLATE_900,
         shape=MSO_SHAPE.RECTANGLE)
add_text(s, Inches(0.8), Inches(2.0), Inches(11.7), Inches(0.7),
         "DS 승강기 유지보수 자재관리 시스템",
         size=18, color=RGBColor(0xBF, 0xDB, 0xFE))
add_text(s, Inches(0.8), Inches(2.7), Inches(11.7), Inches(1.2),
         "자재 신청부터 완료까지",
         size=44, bold=True, color=WHITE)
add_text(s, Inches(0.8), Inches(4.0), Inches(11.7), Inches(0.6),
         "Part 1. 일반 자재 신청 흐름   |   Part 2. 견적 기반 청구 흐름",
         size=16, color=RGBColor(0x93, 0xC5, 0xFD))
add_text(s, Inches(0.8), Inches(5.6), Inches(11.7), Inches(0.5),
         "현장 보수원 (모바일)  →  본사 담당자 (웹)  →  거래처/고객",
         size=13, color=SLATE_200)
add_text(s, Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.4),
         "Next.js 16  ·  Supabase  ·  Phase 1~6 통합 워크플로우",
         size=11, color=SLATE_500)

# ─── 2. 전체 흐름 개요 ───
s = new_slide()
header_bar(s, "전체 흐름 개요", "두 갈래 흐름 — 일반 자재 신청 vs 견적 기반 청구", accent=NAVY)

# Part 1: 일반 흐름
add_text(s, Inches(0.5), Inches(1.15), Inches(12), Inches(0.4),
         "📦 Part 1. 일반 자재 신청 흐름 — 모바일 기사 ↔ 본사 담당자",
         size=14, bold=True, color=BLUE)
labels_p1 = ["자재신청", "확인", "출고", "수령", "회수"]
left = Inches(0.6)
for i, lab in enumerate(labels_p1):
    step_box(s, left, Inches(1.7), Inches(2.1), Inches(1.0), i+1, lab, "", accent=BLUE)
    left += Inches(2.3)
    if i < len(labels_p1) - 1:
        arrow_right(s, left - Inches(0.3), Inches(2.05), color=BLUE)

# Part 2: 견적 흐름
add_text(s, Inches(0.5), Inches(3.5), Inches(12), Inches(0.4),
         "📋 Part 2. 견적 기반 청구 흐름 — 유상 견적 → 청구 → 입금",
         size=14, bold=True, color=AMBER)
labels_p2 = ["청구등록", "견적요청", "견적서", "자재신청", "출고", "세금계산서", "입금"]
left = Inches(0.6)
box_w = Inches(1.55)
for i, lab in enumerate(labels_p2):
    step_box(s, left, Inches(4.05), box_w, Inches(1.0), i+1, lab, "", accent=AMBER)
    left += box_w + Inches(0.18)
    if i < len(labels_p2) - 1:
        arrow_right(s, left - Inches(0.22), Inches(4.4), color=AMBER, width=Inches(0.2))

# 하단 핵심 차이
add_rect(s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.4),
         fill=SLATE_50, line=SLATE_200)
add_text(s, Inches(0.7), Inches(5.65), Inches(12), Inches(0.35),
         "핵심 차이",
         size=12, bold=True, color=SLATE_900)
add_text(s, Inches(0.7), Inches(6.0), Inches(6.0), Inches(0.85),
         "• Part 1: FM(무상) 현장 또는 자체 자재 — 견적 없이 자재만 이동\n"
         "• 청구자=현장 기사 / 처리자=자재 담당자",
         size=11, color=SLATE_700)
add_text(s, Inches(6.9), Inches(6.0), Inches(6.0), Inches(0.85),
         "• Part 2: 외부 청구·유상 작업 — 견적서·세금계산서·입금까지\n"
         "• 청구자=보수원 / 발행자=관리자(admin)",
         size=11, color=SLATE_700)
footer(s, 2, TOTAL_SLIDES)

# ─── 3. Part 1 섹션 표지 ───
s = new_slide()
add_rect(s, Inches(0), Inches(0), WIDE_W, WIDE_H, fill=BLUE, line=BLUE,
         shape=MSO_SHAPE.RECTANGLE)
add_text(s, Inches(0.8), Inches(2.5), Inches(11.7), Inches(0.7),
         "PART 1",
         size=20, color=RGBColor(0xBF, 0xDB, 0xFE))
add_text(s, Inches(0.8), Inches(3.1), Inches(11.7), Inches(1.4),
         "일반 자재 신청 흐름",
         size=48, bold=True, color=WHITE)
add_text(s, Inches(0.8), Inches(4.7), Inches(11.7), Inches(0.5),
         "모바일 신청 → 담당자 확인 → 재고 출고/발주 → 입고 → 수령 → 회수",
         size=14, color=RGBColor(0xDB, 0xEA, 0xFE))
footer(s, 3, TOTAL_SLIDES)

# ─── 4. [1-1] 자재 신청 ───
s = new_slide()
header_bar(s, "1-1. 자재 신청 (현장 기사)", "모바일에서 호기 단위로 자재 신청. 승인 절차 없음.", accent=BLUE)
# 좌측 화면 정보
add_rect(s, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5), fill=SLATE_50, line=SLATE_200)
add_text(s, Inches(0.7), Inches(1.35), Inches(5.5), Inches(0.35),
         "📱 모바일 화면 — /requests/new (예정)",
         size=12, bold=True, color=SLATE_900)
add_text(s, Inches(0.7), Inches(1.8), Inches(5.5), Inches(4.7),
         "입력 항목\n"
         "  • 현장명 + 호기 (자동완성)\n"
         "  • 자재명 또는 자재코드 검색\n"
         "  • 수량 / 비고\n"
         "  • 시급도(긴급/일반)\n\n"
         "특징\n"
         "  • 장갑 환경 고려 — 입력 필드 ≥ 48dp\n"
         "  • 별도 승인 절차 없이 즉시 신청\n"
         "  • 동일 호기 다건 신청 가능\n\n"
         "데이터\n"
         "  • material_requests INSERT\n"
         "  • request_type = 무상 / 당직 / 유상견적",
         size=11, color=SLATE_700)
# 우측 데이터 흐름
add_rect(s, Inches(6.6), Inches(1.2), Inches(6.2), Inches(5.5), fill=WHITE, line=SLATE_200)
add_text(s, Inches(6.8), Inches(1.35), Inches(5.9), Inches(0.35),
         "🗄 핵심 테이블",
         size=12, bold=True, color=SLATE_900)
add_text(s, Inches(6.8), Inches(1.85), Inches(5.9), Inches(4.7),
         "material_requests\n"
         "  ├ request_no  PK\n"
         "  ├ site_name / elevator_name\n"
         "  ├ requester_id / requester_name\n"
         "  ├ status: 신청 → 처리중 → 완료 / 취소\n"
         "  ├ request_type: 무상 / 당직 / 유상견적\n"
         "  └ items[]  (자재코드, 수량, 호기, 비고)\n\n"
         "관련 페이지 (웹)\n"
         "  • /claim/new — 보수원이 PC에서 등록할 때 사용\n"
         "    (유상/무상/당직 3모드 토글 — Phase 3)\n\n"
         "초기 상태\n"
         "  status='신청'",
         size=11, color=SLATE_700)
footer(s, 4, TOTAL_SLIDES)

# ─── 5. [1-2] 신청 확인 ───
s = new_slide()
header_bar(s, "1-2. 신청 확인 (자재 담당자)", "/requests 페이지에서 신규 신청 검토 — 재고 여부 판단", accent=BLUE)
add_rect(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.7), fill=SLATE_50, line=SLATE_200)
add_text(s, Inches(0.8), Inches(1.45), Inches(12), Inches(0.4),
         "💻 웹 화면 — /requests (자재 신청 관리)",
         size=13, bold=True, color=SLATE_900)
# 좌 — 작업
add_text(s, Inches(0.9), Inches(2.0), Inches(5.7), Inches(0.4),
         "담당자 작업",
         size=12, bold=True, color=BLUE)
add_text(s, Inches(0.9), Inches(2.4), Inches(5.7), Inches(4.3),
         "✓ 신청 내역 검토\n"
         "    현장·호기·자재·수량·신청자 일괄 조회\n\n"
         "✓ 재고 현황 확인\n"
         "    자재별 가용 재고 옆에 표시\n\n"
         "✓ 분기 처리\n"
         "    재고 있음  →  바로 출고\n"
         "    재고 없음  →  발주 등록\n\n"
         "✓ 상태 갱신\n"
         "    status: 신청 → 처리중",
         size=11, color=SLATE_700)
# 우 — 의사결정 다이어그램
left = Inches(7.2)
add_text(s, left, Inches(2.0), Inches(5.5), Inches(0.4),
         "의사결정",
         size=12, bold=True, color=BLUE)
step_box(s, left, Inches(2.5), Inches(2.4), Inches(0.85),
         "?", "재고 확인", "현재 가용 수량 ≥ 신청 수량?", accent=SLATE_500)
arrow_down(s, left + Inches(1.05), Inches(3.4))
# 분기 두 박스
step_box(s, left, Inches(3.85), Inches(2.4), Inches(0.85),
         "Y", "재고 출고", "/outbound 즉시 출고 처리", accent=EMERALD)
step_box(s, left + Inches(2.7), Inches(3.85), Inches(2.4), Inches(0.85),
         "N", "발주 등록", "/purchase-orders 발주 생성", accent=AMBER)
footer(s, 5, TOTAL_SLIDES)

# ─── 6. [1-3] 재고 출고 ───
s = new_slide()
header_bar(s, "1-3. 재고 출고 (재고 보유 시)", "재고가 있으면 즉시 출고 처리 — add_transaction RPC로 재고 차감", accent=EMERALD)
add_text(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.4),
         "💻 /outbound — 출고 관리",
         size=13, bold=True, color=SLATE_900)
# 4개 단계 박스
steps = [
    ("출고 등록", "신청 행을 선택 후 [출고]\nS/N 자재는 시리얼 지정"),
    ("재고 차감", "add_transaction RPC\ntype='출고', qty, site, elevator"),
    ("trans 기록", "transactions 행 생성\n· user_name / created_at"),
    ("상태 갱신", "material_requests\nstatus='완료'\n신청자에게 카카오 알림"),
]
left = Inches(0.6)
box_w = Inches(2.9)
for i, (title, body) in enumerate(steps):
    step_box(s, left, Inches(1.9), box_w, Inches(1.5), i+1, title, body, accent=EMERALD)
    left += box_w + Inches(0.18)
    if i < len(steps) - 1:
        arrow_right(s, left - Inches(0.22), Inches(2.5), color=EMERALD, width=Inches(0.2))

# 데이터 효과
add_rect(s, Inches(0.5), Inches(3.7), Inches(12.3), Inches(3.4), fill=SLATE_50, line=SLATE_200)
add_text(s, Inches(0.8), Inches(3.85), Inches(12), Inches(0.4),
         "데이터 효과",
         size=12, bold=True, color=SLATE_900)
add_text(s, Inches(0.8), Inches(4.3), Inches(6.0), Inches(2.7),
         "transactions 새 행\n"
         "  type='출고', material_id, qty\n"
         "  site_name, elevator_name\n"
         "  user_id (담당자) / user_name\n"
         "  created_at = 출고 일시\n\n"
         "재고 변동 (실시간 합산)\n"
         "  자재별 잔여 수량 = Σ(입고 − 출고 − 반납회수 등)",
         size=11, color=SLATE_700)
add_text(s, Inches(6.9), Inches(4.3), Inches(6.0), Inches(2.7),
         "S/N 추적 자재\n"
         "  • 시리얼별 instance 상태 = '출고됨'\n"
         "  • /serial-history 에서 타임라인 조회\n\n"
         "회수 표시\n"
         "  • 출고 시 회수 필요 자재면 needs_return=true\n"
         "  • 추후 /returns 에서 반납 등록",
         size=11, color=SLATE_700)
footer(s, 6, TOTAL_SLIDES)

# ─── 7. [1-4] 발주 (재고 없음) ───
s = new_slide()
header_bar(s, "1-4. 발주 (재고 부족 시)", "/purchase-orders — 거래처에 발주 등록, 입고 대기", accent=AMBER)
add_rect(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.7), fill=SLATE_50, line=SLATE_200)

# 발주 흐름
left = Inches(0.8)
add_text(s, left, Inches(1.4), Inches(12), Inches(0.4),
         "발주 흐름",
         size=12, bold=True, color=AMBER)
add_text(s, left, Inches(1.85), Inches(12), Inches(2.0),
         "① 자재 담당자가 거래처 선택 + 발주 라인 입력\n"
         "    (자재코드, 수량, 단가, 납기일)\n\n"
         "② purchase_orders + purchase_order_items INSERT\n"
         "    order_no: PO-YYYY-NNNN 자동 채번\n\n"
         "③ 상태: 발주 → 부분입고 → 입고완료 → 종료",
         size=11, color=SLATE_700)

# 발주 ↔ 신청 연결
add_text(s, left, Inches(4.1), Inches(12), Inches(0.4),
         "신청과의 연결",
         size=12, bold=True, color=AMBER)
add_text(s, left, Inches(4.55), Inches(12), Inches(2.3),
         "• material_requests.status = '처리중' 유지\n"
         "• 발주 라인이 입고 완료 시 자동으로 출고 가능 상태로 전환\n"
         "• 동일 자재 다중 신청은 발주 1건으로 묶어서 처리 가능\n"
         "• 거래처별 발주 이력은 /data/vendors 거래처 상세에서도 조회",
         size=11, color=SLATE_700)
footer(s, 7, TOTAL_SLIDES)

# ─── 8. [1-5] 입고 처리 ───
s = new_slide()
header_bar(s, "1-5. 입고 처리 (발주 자재 도착)", "/inbound — 발주된 자재 입고 등록 + 신청 기사에게 카카오 알림", accent=EMERALD)

# 좌 : 입고 단계
add_text(s, Inches(0.5), Inches(1.2), Inches(6.2), Inches(0.4),
         "입고 처리 단계",
         size=12, bold=True, color=EMERALD)
items = [
    "발주 라인 선택",
    "실제 입고 수량 + S/N 입력",
    "transactions 'type=입고' 행 생성",
    "purchase_order_items.received_qty 누적",
    "전 라인 완료 시 PO.status='입고완료'",
]
top = Inches(1.7)
for i, t in enumerate(items):
    step_box(s, Inches(0.5), top, Inches(6.0), Inches(0.7), i+1, t, "", accent=EMERALD)
    top += Inches(0.85)

# 우 : 카카오 알림
add_rect(s, Inches(7.0), Inches(1.2), Inches(5.8), Inches(5.7), fill=SLATE_50, line=SLATE_200)
add_text(s, Inches(7.2), Inches(1.35), Inches(5.5), Inches(0.4),
         "📨 카카오 알림톡",
         size=12, bold=True, color=ROSE)
add_text(s, Inches(7.2), Inches(1.8), Inches(5.5), Inches(5.0),
         "발송 시점\n"
         "  • 신청한 자재가 입고 완료된 직후\n\n"
         "수신 대상\n"
         "  • material_requests.requester_id\n"
         "  • users.phone 으로 발송\n\n"
         "메시지 예시\n"
         "  [DS 자재 알림]\n"
         "  요청하신 자재가 입고되었습니다.\n"
         "  · 현장: 청담아이파크 #3호기\n"
         "  · 자재: VVVF 인버터 1EA\n"
         "  · 수령처: 본사 자재창고\n\n"
         "API\n"
         "  알리고 / 비즈톡 등 (현재 미연동)",
         size=11, color=SLATE_700)
footer(s, 8, TOTAL_SLIDES)

# ─── 9. [1-6] 수령 완료 ───
s = new_slide()
header_bar(s, "1-6. 수령 완료 (기사 최종 수령)", "기사가 자재를 받은 시점 — 신청 종결", accent=EMERALD)

add_rect(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.7), fill=WHITE, line=SLATE_200)

# 좌 — 수령 처리
add_text(s, Inches(0.8), Inches(1.4), Inches(5.7), Inches(0.4),
         "수령 등록 (모바일)",
         size=12, bold=True, color=EMERALD)
add_text(s, Inches(0.8), Inches(1.85), Inches(5.7), Inches(4.9),
         "• 기사가 자재를 인수한 뒤 모바일에서 [수령확인]\n\n"
         "• 사진 첨부(옵션) — 자재 수령 인증\n\n"
         "• material_requests.status = '완료'\n"
         "    completed_at = 현재 일시 기록\n\n"
         "• 신청자에 '완료' 푸시 알림\n\n"
         "• 출고가 본사 직접 픽업이면\n"
         "    1-3 출고 시점에 바로 '완료' 처리",
         size=11, color=SLATE_700)

# 우 — 상태 전이 다이어그램
add_text(s, Inches(7.0), Inches(1.4), Inches(5.7), Inches(0.4),
         "신청 상태 전이",
         size=12, bold=True, color=EMERALD)
states = [
    ("신청",     SLATE_500),
    ("처리중",   AMBER),
    ("출고완료", BLUE),
    ("수령완료", EMERALD),
]
top = Inches(1.85)
for i, (state, color) in enumerate(states):
    step_box(s, Inches(7.0), top, Inches(5.7), Inches(0.7),
             i+1, state, "", accent=color)
    top += Inches(0.95)
    if i < len(states) - 1:
        arrow_down(s, Inches(9.7), top - Inches(0.32))
footer(s, 9, TOTAL_SLIDES)

# ─── 10. [1-7] 회수/반납 ───
s = new_slide()
header_bar(s, "1-7. 회수 / 반납 (교체 후 옛 부품 회수)", "/returns — 출고 시 회수 표시된 자재는 반납 등록까지 추적", accent=ROSE)

add_rect(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.7), fill=SLATE_50, line=SLATE_200)

# 좌 — 트리거
add_text(s, Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.4),
         "회수 등록 시점",
         size=12, bold=True, color=ROSE)
add_text(s, Inches(0.8), Inches(1.85), Inches(5.5), Inches(4.9),
         "교체 작업이 끝나면 기사가 모바일에서 회수 등록\n\n"
         "필수 입력\n"
         "  • 자재 (출고 이력에서 선택)\n"
         "  • 시리얼번호 (S/N 자재의 경우)\n"
         "  • 회수 상태: 양호 / 손상 / 폐기\n\n"
         "회수 사진\n"
         "  • Storage(material-returns) 업로드\n\n"
         "transactions 'type=회수' 행 생성\n"
         "  → S/N instance.status='회수됨'\n"
         "  → 본사 입고 검수 후 '재사용 수리품(R)' 분류 가능",
         size=11, color=SLATE_700)

# 우 — 후속 처리
add_text(s, Inches(7.0), Inches(1.4), Inches(5.7), Inches(0.4),
         "수리품 재투입 (옵션)",
         size=12, bold=True, color=ROSE)
add_text(s, Inches(7.0), Inches(1.85), Inches(5.7), Inches(4.9),
         "자재코드 12자리 중\n"
         "  마지막 1자리 '수리품 구분' = R\n"
         "  → 같은 부품의 재생품으로 분류\n\n"
         "수리품으로 등록되면\n"
         "  • materials.warranty_months = 12 (자동)\n"
         "  • 신품(_)과 재고 분리\n"
         "  • 출고 시 [수리품/신품] 선택 가능\n\n"
         "/serial-history\n"
         "  S/N 단위 입고→출고→회수→재출고\n"
         "  타임라인을 한 화면에서 확인",
         size=11, color=SLATE_700)
footer(s, 10, TOTAL_SLIDES)

# ─── 11. Part 2 섹션 표지 ───
s = new_slide()
add_rect(s, Inches(0), Inches(0), WIDE_W, WIDE_H, fill=AMBER, line=AMBER,
         shape=MSO_SHAPE.RECTANGLE)
add_text(s, Inches(0.8), Inches(2.5), Inches(11.7), Inches(0.7),
         "PART 2",
         size=20, color=RGBColor(0xFE, 0xF3, 0xC7))
add_text(s, Inches(0.8), Inches(3.1), Inches(11.7), Inches(1.4),
         "견적 기반 청구 흐름",
         size=48, bold=True, color=WHITE)
add_text(s, Inches(0.8), Inches(4.7), Inches(11.7), Inches(0.5),
         "청구 → 견적서 → 자재출고 → 거래명세서·세금계산서 → 입금 → 종료",
         size=14, color=RGBColor(0xFE, 0xF3, 0xC7))
add_text(s, Inches(0.8), Inches(5.4), Inches(11.7), Inches(0.4),
         "Phase 1~6 통합 워크플로우 (2026-05-16 완성)",
         size=12, color=RGBColor(0xFE, 0xF3, 0xC7))
footer(s, 11, TOTAL_SLIDES)

# ─── 12. [2-1] 보수원 청구 등록 ───
s = new_slide()
header_bar(s, "2-1. 보수원 청구 등록 (/claim/new)", "유상 / 무상 / 당직선출고 3 모드 토글 — 한 화면에서 분기", accent=AMBER)

# 3 모드 카드
modes = [
    ("유상견적요청", "외부 청구·유상 작업\n→ quote_requests INSERT\n→ 견적서 작성 트리거", AMBER),
    ("무상신청",     "FM(무상) 현장 또는\n자체 자재 무상 처리\n→ material_requests INSERT", EMERALD),
    ("당직선출고",   "야간·당직 시 선출고\n→ material_requests (pending)\n→ 추후 견적 사후 처리", ROSE),
]
left = Inches(0.5)
for title, body, color in modes:
    step_box(s, left, Inches(1.3), Inches(4.1), Inches(2.5), "•", title, body, accent=color)
    left += Inches(4.3)

# 자동 분기
add_rect(s, Inches(0.5), Inches(4.05), Inches(12.3), Inches(3.0), fill=SLATE_50, line=SLATE_200)
add_text(s, Inches(0.8), Inches(4.2), Inches(12), Inches(0.4),
         "자동 분기 규칙",
         size=12, bold=True, color=SLATE_900)
add_text(s, Inches(0.8), Inches(4.65), Inches(12), Inches(2.3),
         "• 현장 선택 시 sites.contract_type = TK-FM / 대솔FM / DS-FM 이면\n"
         "   자동으로 [무상] 모드로 토글 + 권유 메시지\n\n"
         "• 자재 라인 입력 — 자재 검색 + 기본 5줄, 자재별 호기 컬럼\n\n"
         "• 등록 후 사이드바 [📝 견적요청 목록] / [📋 자재 신청 관리] 에서 추적",
         size=11, color=SLATE_700)
footer(s, 12, TOTAL_SLIDES)

# ─── 13. [2-2] 견적요청 → 견적서 ───
s = new_slide()
header_bar(s, "2-2. 견적요청 → 견적서 작성", "/claim/quote-requests 목록 → [📝 견적서 작성] → /quotes/new?fromRequest=N", accent=AMBER)

# 단계 박스
steps = [
    ("요청 검토", "좌목록·우상세 2단\n자재·호기·요청자 검토"),
    ("견적 진입", "[📝 견적서 작성] 클릭\n?fromRequest=N 으로 라우팅"),
    ("자동 프리필", "헤더·items 자동 입력\nstatus='견적작성중'"),
    ("견적 발행", "자재 단가 보강, 인건비 입력\n저장 시 quote_id 연결\nstatus='견적발행'"),
]
left = Inches(0.5)
box_w = Inches(2.95)
for i, (title, body) in enumerate(steps):
    step_box(s, left, Inches(1.3), box_w, Inches(1.7), i+1, title, body, accent=AMBER)
    left += box_w + Inches(0.18)
    if i < len(steps) - 1:
        arrow_right(s, left - Inches(0.22), Inches(2.0), color=AMBER, width=Inches(0.2))

# 데이터 연결
add_rect(s, Inches(0.5), Inches(3.3), Inches(12.3), Inches(3.7), fill=SLATE_50, line=SLATE_200)
add_text(s, Inches(0.8), Inches(3.45), Inches(12), Inches(0.4),
         "데이터 연결",
         size=12, bold=True, color=SLATE_900)
add_text(s, Inches(0.8), Inches(3.9), Inches(12), Inches(3.0),
         "quote_requests\n"
         "  ├ request_no  (QR-YYYY-NNNN)\n"
         "  ├ status: 신청 → 견적작성중 → 견적발행 → 취소\n"
         "  └ quote_id  ←  FK to quotes.id  (발행 시 자동 연결)\n\n"
         "견적서 헤더에 '📋 견적요청 QR-XXXX-NNNN 에서 가져옴' 배지 표시\n"
         "→ 양방향 추적 가능 (요청 ↔ 견적서)",
         size=11, color=SLATE_700)
footer(s, 13, TOTAL_SLIDES)

# ─── 14. [2-3] 견적서 결재 + 진행상태 ───
s = new_slide()
header_bar(s, "2-3. 견적서 결재 + 진행상태 스텝퍼", "결재(작성중→발행→승인) 와 진행상태(미시작→...→종료) 가 분리되어 운영", accent=AMBER)

# 결재
add_text(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.4),
         "결재 상태 (status)",
         size=12, bold=True, color=BLUE)
status_flow = [("작성중", SLATE_500), ("발행", BLUE), ("승인", EMERALD), ("취소", ROSE)]
left = Inches(0.7)
for i, (st, color) in enumerate(status_flow):
    step_box(s, left, Inches(1.65), Inches(2.5), Inches(0.85),
             "✓" if i < 3 else "✕", st, "", accent=color)
    left += Inches(2.7)
    if i < 2:
        arrow_right(s, left - Inches(0.22), Inches(2.0), color=SLATE_500, width=Inches(0.2))

# 진행상태
add_text(s, Inches(0.5), Inches(2.85), Inches(12.3), Inches(0.4),
         "진행상태 (progress_state) — 승인 후에만 이동 가능",
         size=12, bold=True, color=AMBER)
progress = ["미시작", "자재신청", "자재출고", "세금계산서발급", "입금완료", "종료"]
left = Inches(0.55)
for i, st in enumerate(progress):
    step_box(s, left, Inches(3.3), Inches(2.0), Inches(0.85), i+1, st, "", accent=AMBER)
    left += Inches(2.13)
    if i < len(progress) - 1:
        arrow_right(s, left - Inches(0.16), Inches(3.65), color=AMBER, width=Inches(0.15))

# 스냅샷
add_rect(s, Inches(0.5), Inches(4.6), Inches(12.3), Inches(2.4), fill=SLATE_50, line=SLATE_200)
add_text(s, Inches(0.8), Inches(4.75), Inches(12), Inches(0.4),
         "📜 quote_revisions  — 모든 상태 전이 직전 자동 스냅샷",
         size=12, bold=True, color=SLATE_900)
add_text(s, Inches(0.8), Inches(5.2), Inches(12), Inches(1.8),
         "• snapshot_quote(quote_id, summary, user_id, user_name)  RPC 호출\n"
         "• revision_no 자동 채번, 헤더·items 전체 스냅샷 보관\n"
         "• [📜 수정 이력] 모달에서 변경 시각·변경자·요약 조회\n"
         "• 작성중 [✏️ 수정] 저장 시에도 동일하게 스냅샷 기록",
         size=11, color=SLATE_700)
footer(s, 14, TOTAL_SLIDES)

# ─── 15. [2-4] 자재신청 자동생성 → 출고 ───
s = new_slide()
header_bar(s, "2-4. 자재신청 자동생성 → 출고 관리", "결재 승인 + 진행상태 미시작 → [📦 자재신청 생성] → /claim/quote-outbound", accent=EMERALD)

# 좌 — 자재신청 자동생성
add_text(s, Inches(0.5), Inches(1.2), Inches(6.0), Inches(0.4),
         "① 자재신청 자동생성",
         size=12, bold=True, color=EMERALD)
add_text(s, Inches(0.5), Inches(1.65), Inches(6.0), Inches(2.7),
         "[📦 자재신청 생성] 버튼 클릭\n\n"
         "→ quote_items 의 각 라인을\n"
         "    material_requests.items 로 자동 매핑\n\n"
         "→ material_requests.quote_id = 견적서 ID\n"
         "    request_type = '유상견적'\n\n"
         "→ 견적서 progress_state = '자재신청' 전이\n"
         "    snapshot_quote 자동 호출\n\n"
         "※ 수기 라인(material_id 미연결)은 제외",
         size=11, color=SLATE_700)

# 우 — 출고 관리
add_text(s, Inches(7.0), Inches(1.2), Inches(5.7), Inches(0.4),
         "② 견적 출고 관리 (/claim/quote-outbound)",
         size=12, bold=True, color=EMERALD)
add_text(s, Inches(7.0), Inches(1.65), Inches(5.7), Inches(2.7),
         "• 라인별 현재 재고 표시\n"
         "• [출고 N] / [발주] 버튼 / [📤 일괄 출고]\n"
         "• S/N 추적 자재는 별도 [출고 관리]로 안내\n"
         "• add_transaction RPC 로 재고 차감 + tx 생성\n\n"
         "✅ 모든 라인 잔여 = 0 → [✅ 출고 완료]\n"
         "  → material_requests.status = '완료'\n"
         "  → quotes.progress_state = '자재출고' 자동",
         size=11, color=SLATE_700)

# 하단 흐름
add_rect(s, Inches(0.5), Inches(4.7), Inches(12.3), Inches(2.3), fill=SLATE_50, line=SLATE_200)
add_text(s, Inches(0.8), Inches(4.85), Inches(12), Inches(0.4),
         "데이터 흐름 요약",
         size=12, bold=True, color=SLATE_900)
add_text(s, Inches(0.8), Inches(5.3), Inches(12), Inches(1.6),
         "quotes → material_requests → transactions(type=출고)\n"
         "          (quote_id FK)            (재고 차감)\n\n"
         "양방향 추적: 견적 상세에서 자재신청 ID 확인 / 자재신청 목록에서 견적번호 표시",
         size=11, color=SLATE_700)
footer(s, 15, TOTAL_SLIDES)

# ─── 16. [2-5] 거래명세서·세금계산서 ───
s = new_slide()
header_bar(s, "2-5. 거래명세서 · 세금계산서 발행", "유상 + 승인 일 때만 발행. 무상 견적은 차단.", accent=BLUE)

# 좌 — 거래명세서
add_rect(s, Inches(0.5), Inches(1.2), Inches(6.1), Inches(5.7), fill=WHITE, line=EMERALD, line_w=1.5)
add_text(s, Inches(0.7), Inches(1.4), Inches(5.7), Inches(0.45),
         "📄 거래명세서  (DN-YYYY-NNNN)",
         size=13, bold=True, color=EMERALD)
add_text(s, Inches(0.7), Inches(1.95), Inches(5.7), Inches(4.8),
         "용도: 자재 출고 시 동봉\n\n"
         "양식 (/invoice/delivery?id=N)\n"
         "  • 공급자·인수자 칸\n"
         "  • 자재 라인 + 합계\n"
         "  • 인수자 서명란\n"
         "  • A4 인쇄 양식\n\n"
         "발행 시\n"
         "  invoices INSERT (invoice_type='거래명세서')\n"
         "  새 탭 자동 오픈 → Ctrl+P 인쇄\n\n"
         "진행상태 영향: 없음",
         size=11, color=SLATE_700)

# 우 — 세금계산서
add_rect(s, Inches(6.7), Inches(1.2), Inches(6.1), Inches(5.7), fill=WHITE, line=BLUE, line_w=1.5)
add_text(s, Inches(6.9), Inches(1.4), Inches(5.7), Inches(0.45),
         "🧾 세금계산서  (TX-YYYY-NNNN)",
         size=13, bold=True, color=BLUE)
add_text(s, Inches(6.9), Inches(1.95), Inches(5.7), Inches(4.8),
         "용도: 회계 처리용\n\n"
         "양식 (/invoice/tax?id=N)\n"
         "  • 공급자 / 공급받는자\n"
         "  • 공급가액·부가세·합계 분리\n"
         "  • 사내 양식 (홈택스/위하고 별도 발급)\n\n"
         "발행 시\n"
         "  invoices INSERT (invoice_type='세금계산서')\n"
         "  vat_amount = supply / 10\n\n"
         "진행상태 영향\n"
         "  progress_state → '세금계산서발급' 자동",
         size=11, color=SLATE_700)
footer(s, 16, TOTAL_SLIDES)

# ─── 17. [2-6] 입금 관리 ───
s = new_slide()
header_bar(s, "2-6. 입금 관리", "1:N 분할 입금 지원 + 누적 도달 시 자동 입금완료", accent=ROSE)

# 좌 — 입금 등록
add_text(s, Inches(0.5), Inches(1.2), Inches(5.7), Inches(0.4),
         "💵 입금 등록 모달",
         size=12, bold=True, color=ROSE)
add_text(s, Inches(0.5), Inches(1.65), Inches(5.7), Inches(5.2),
         "입력 항목\n"
         "  • 금액 (필수)\n"
         "  • 방법: 현금 / 계좌이체 / 카드 / 어음 / 기타\n"
         "  • 입금일 / 입금자명\n"
         "  • 선입금 플래그 (is_prepaid)\n"
         "  • 비고\n\n"
         "데이터\n"
         "  payments INSERT\n"
         "  payment_no: PAY-YYYY-NNNN\n"
         "  status='확정' (기본)\n\n"
         "자동 전이\n"
         "  Σ확정입금 ≥ total_amount\n"
         "  → progress_state='입금완료'",
         size=11, color=SLATE_700)

# 우 — 입금 이력
add_text(s, Inches(7.0), Inches(1.2), Inches(5.7), Inches(0.4),
         "💼 입금 이력 모달",
         size=12, bold=True, color=ROSE)
add_text(s, Inches(7.0), Inches(1.65), Inches(5.7), Inches(5.2),
         "표시 정보\n"
         "  • 입금번호 / 일자 / 금액 / 방법\n"
         "  • 입금자명 / 등록자 / 비고\n"
         "  • 선입금 뱃지\n"
         "  • 취소 처리 가능\n\n"
         "요약\n"
         "  • 총액 ₩XXX\n"
         "  • 누적입금 ₩YYY\n"
         "  • 미수금 ₩(총액−누적)\n"
         "  • 완납 시 'PAID' 강조\n\n"
         "→ /stats/performance 사원 실적에 반영",
         size=11, color=SLATE_700)
footer(s, 17, TOTAL_SLIDES)

# ─── 18. [2-7] 견적 종료 + 마무리 ───
s = new_slide()
header_bar(s, "2-7. 견적 종료 + 워크플로우 마무리", "[🏁 견적 종료] 클릭 시 progress_state='종료' — 전 흐름 완료", accent=SLATE_900)

# 진입 조건
add_text(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.4),
         "🏁 견적 종료 조건",
         size=12, bold=True, color=SLATE_900)
add_text(s, Inches(0.5), Inches(1.7), Inches(12.3), Inches(0.8),
         "결재 = 승인  +  charge_type = 유상  +  progress_state ∈ {자재출고, 세금계산서발급, 입금완료}",
         size=11, color=SLATE_700)

# 후속 영향
add_rect(s, Inches(0.5), Inches(2.6), Inches(6.0), Inches(4.4), fill=SLATE_50, line=SLATE_200)
add_text(s, Inches(0.7), Inches(2.75), Inches(5.7), Inches(0.4),
         "데이터 효과",
         size=12, bold=True, color=SLATE_900)
add_text(s, Inches(0.7), Inches(3.2), Inches(5.7), Inches(3.7),
         "• progress_state = '종료' 갱신\n"
         "• snapshot_quote 자동 호출\n"
         "    summary='진행상태: ... → 종료'\n"
         "• 견적 상세 헤더에 'CLOSED' 표시\n\n"
         "• 견적은 읽기 전용으로 잠김\n"
         "    (admin 도 수정 불가)\n\n"
         "• 미수금이 있어도 종료 가능\n"
         "    (사후 추가 입금 시 [확정] 처리)",
         size=11, color=SLATE_700)

# 사원 실적 반영
add_rect(s, Inches(6.7), Inches(2.6), Inches(6.1), Inches(4.4), fill=SLATE_50, line=SLATE_200)
add_text(s, Inches(6.9), Inches(2.75), Inches(5.7), Inches(0.4),
         "📈 사원 실적 (매출) 반영",
         size=12, bold=True, color=SLATE_900)
add_text(s, Inches(6.9), Inches(3.2), Inches(5.7), Inches(3.7),
         "/stats/performance — 월/분기/년 토글\n\n"
         "포함 범위\n"
         "  • 매출확정 — 자재출고 이상 진행상태\n"
         "  • 전체 — 취소 제외 모두\n\n"
         "사원별 정렬\n"
         "  🥇 🥈 🥉 + 비중바 시각화\n"
         "  매출 / 입금 / 미수금 칼럼\n\n"
         "→ Phase 6 종합 대시보드로 흐름 종결",
         size=11, color=SLATE_700)
footer(s, 18, TOTAL_SLIDES)

# ──────────────────────────── 저장 ────────────────────────────
out_dir = Path("docs")
out_dir.mkdir(exist_ok=True)
out_path = out_dir / "자재흐름_DS_Material.pptx"
prs.save(out_path)
print(f"[OK] Saved: {out_path.resolve()}")
print(f"     slides: {len(prs.slides)}")
