"""
DS_Material — 견적요청 → 세금계산서 발행 임직원 설명용 PPT 생성
실행: python scripts/generate_invoice_flow_ppt.py
산출: docs/견적-세금계산서_발행안내_DS_Material.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

# ── 색상 ──
NAVY      = RGBColor(0x1E, 0x3A, 0x8A)
BLUE      = RGBColor(0x25, 0x63, 0xEB)
AMBER     = RGBColor(0xD9, 0x77, 0x06)
EMERALD   = RGBColor(0x05, 0x96, 0x69)
ROSE      = RGBColor(0xE1, 0x1D, 0x48)
VIOLET    = RGBColor(0x7C, 0x3A, 0xED)
SLATE_900 = RGBColor(0x0F, 0x17, 0x2A)
SLATE_700 = RGBColor(0x33, 0x41, 0x55)
SLATE_500 = RGBColor(0x64, 0x74, 0x8B)
SLATE_200 = RGBColor(0xE2, 0xE8, 0xF0)
SLATE_100 = RGBColor(0xF1, 0xF5, 0xF9)
SLATE_50  = RGBColor(0xF8, 0xFA, 0xFC)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
RED_50    = RGBColor(0xFE, 0xF2, 0xF2)
AMBER_50  = RGBColor(0xFF, 0xFB, 0xEB)
BLUE_50   = RGBColor(0xEF, 0xF6, 0xFF)
EMER_50   = RGBColor(0xEC, 0xFD, 0xF5)

WIDE_W, WIDE_H = Inches(13.333), Inches(7.5)


def add_text(slide, left, top, width, height, text, *,
             size=14, bold=False, color=SLATE_900, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, font="맑은 고딕"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_rect(slide, left, top, width, height, *,
             fill=WHITE, line=SLATE_200, line_w=0.75, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    s = slide.shapes.add_shape(shape, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s


def arrow_right(slide, left, top, width=Inches(0.4), height=Inches(0.3), color=SLATE_500):
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s


def arrow_down(slide, left, top, width=Inches(0.3), height=Inches(0.4), color=SLATE_500):
    s = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s


def header_bar(slide, title, subtitle, *, accent=BLUE):
    add_rect(slide, Inches(0), Inches(0), WIDE_W, Inches(0.9), fill=SLATE_50, line=None, shape=MSO_SHAPE.RECTANGLE)
    add_rect(slide, Inches(0), Inches(0), Inches(0.18), Inches(0.9), fill=accent, line=None, shape=MSO_SHAPE.RECTANGLE)
    add_text(slide, Inches(0.4), Inches(0.12), Inches(12.5), Inches(0.45), title, size=22, bold=True, color=SLATE_900)
    add_text(slide, Inches(0.4), Inches(0.54), Inches(12.5), Inches(0.32), subtitle, size=12, color=SLATE_500)


def footer(slide, n, total):
    add_text(slide, Inches(0.4), Inches(7.12), Inches(12), Inches(0.3),
             f"DS_Material — 견적요청 ▸ 세금계산서 발행 안내   ·   {n} / {total}",
             size=9, color=SLATE_500)


def badge(slide, left, top, num, accent):
    sz = Inches(0.5)
    add_rect(slide, left, top, sz, sz, fill=accent, line=None, shape=MSO_SHAPE.OVAL)
    add_text(slide, left, top, sz, sz, str(num), size=15, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def step_detail(slide, num, accent, who, where, what, result, *, top=Inches(1.25)):
    """단계 상세: 좌측 큰 번호 + 4개 정보 카드(담당/화면/하는 일/결과)"""
    # 좌측 번호 패널
    add_rect(slide, Inches(0.5), top, Inches(2.3), Inches(5.3), fill=accent, line=None)
    add_text(slide, Inches(0.5), top + Inches(1.4), Inches(2.3), Inches(1.3), str(num),
             size=110, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(0.5), top + Inches(3.7), Inches(2.3), Inches(1.2), "STEP",
             size=18, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)

    cards = [
        ("👤  담당", who, BLUE_50, BLUE),
        ("🖥  화면 / 메뉴", where, SLATE_50, SLATE_500),
        ("✍  하는 일", what, AMBER_50, AMBER),
        ("✅  시스템 결과", result, EMER_50, EMERALD),
    ]
    cx = Inches(3.05)
    cw = Inches(9.75)
    ch = Inches(1.2)
    gap = Inches(0.12)
    cy = top
    for label, body, fill, ac in cards:
        add_rect(slide, cx, cy, cw, ch, fill=fill, line=SLATE_200, line_w=0.75)
        add_rect(slide, cx, cy, Inches(0.08), ch, fill=ac, line=None, shape=MSO_SHAPE.RECTANGLE)
        add_text(slide, cx + Inches(0.25), cy + Inches(0.1), Inches(3), Inches(0.35),
                 label, size=11, bold=True, color=ac if ac != SLATE_500 else SLATE_700)
        add_text(slide, cx + Inches(0.25), cy + Inches(0.46), cw - Inches(0.5), ch - Inches(0.5),
                 body, size=12, color=SLATE_900)
        cy = cy + ch + gap


prs = Presentation()
prs.slide_width = WIDE_W
prs.slide_height = WIDE_H
BLANK = prs.slide_layouts[6]
TOTAL = 12


def new_slide():
    return prs.slides.add_slide(BLANK)


# ─── 1. 표지 ───
s = new_slide()
add_rect(s, Inches(0), Inches(0), WIDE_W, WIDE_H, fill=NAVY, line=None, shape=MSO_SHAPE.RECTANGLE)
add_rect(s, Inches(0), Inches(5.0), WIDE_W, Inches(2.5), fill=SLATE_900, line=None, shape=MSO_SHAPE.RECTANGLE)
add_text(s, Inches(0.8), Inches(1.7), Inches(11.7), Inches(0.6),
         "DS 승강기 유지보수 자재관리 시스템", size=18, color=RGBColor(0xBF, 0xDB, 0xFE))
add_text(s, Inches(0.8), Inches(2.4), Inches(11.7), Inches(1.3),
         "견적요청부터 세금계산서 발행까지", size=42, bold=True, color=WHITE)
add_text(s, Inches(0.8), Inches(3.8), Inches(11.7), Inches(0.6),
         "임직원 업무 안내   ·   6단계로 보는 유상 청구 흐름", size=16, color=RGBColor(0x93, 0xC5, 0xFD))
add_text(s, Inches(0.8), Inches(5.4), Inches(11.7), Inches(0.5),
         "현장 보수원  →  견적 담당  →  자재 담당  →  회계/세무", size=14, color=SLATE_200)
add_text(s, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.4),
         "견적요청 ▸ 견적서 ▸ 결재승인 ▸ 자재출고 ▸ 세금계산서 발행 ▸ 입금", size=12, color=SLATE_500)

# ─── 2. 전체 흐름 한눈에 ───
s = new_slide()
header_bar(s, "전체 흐름 한눈에", "유상 자재 청구는 아래 6단계를 거칩니다", accent=NAVY)
steps = [
    ("1", "견적요청", "현장 보수원", BLUE),
    ("2", "견적서 작성", "견적 담당", AMBER),
    ("3", "결재 승인", "관리자", VIOLET),
    ("4", "자재 신청·출고", "자재 담당", EMERALD),
    ("5", "세금계산서 발행", "회계/담당", ROSE),
    ("6", "입금 관리", "회계", NAVY),
]
n = len(steps)
bw, bh = Inches(1.78), Inches(2.2)
gap = Inches(0.18)
total_w = bw * n + gap * (n - 1)
x0 = (WIDE_W - total_w) / 2
y0 = Inches(2.3)
for i, (num, title, who, ac) in enumerate(steps):
    x = x0 + i * (bw + gap)
    add_rect(s, x, y0, bw, bh, fill=WHITE, line=SLATE_200, line_w=1.0)
    add_rect(s, x, y0, bw, Inches(0.12), fill=ac, line=None, shape=MSO_SHAPE.RECTANGLE)
    badge(s, x + (bw - Inches(0.5)) / 2, y0 + Inches(0.35), num, ac)
    add_text(s, x, y0 + Inches(1.0), bw, Inches(0.6), title, size=13, bold=True,
             color=SLATE_900, align=PP_ALIGN.CENTER)
    add_text(s, x, y0 + Inches(1.6), bw, Inches(0.5), who, size=10, color=SLATE_500,
             align=PP_ALIGN.CENTER)
    if i < n - 1:
        arrow_right(s, x + bw + Inches(0.0), y0 + bh / 2 - Inches(0.15),
                    width=Inches(0.18), height=Inches(0.3), color=SLATE_500)
add_text(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.5),
         "※ 무상(FM 현장) 자재는 견적·발행 대상이 아니며 '자재신청'으로 처리합니다.",
         size=12, bold=True, color=AMBER, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(5.6), Inches(12.3), Inches(1.2),
         "각 단계는 시스템에 자동 기록되며, 견적서 한 건을 기준으로 신청 → 출고 → 발행 → 입금이 모두 연결되어 추적됩니다.",
         size=12, color=SLATE_700, align=PP_ALIGN.CENTER)
footer(s, 2, TOTAL)

# ─── 3~8. 단계별 상세 ───
details = [
    (1, BLUE, "현장 보수원 (모바일/웹)",
     "메뉴: [견적 및 자재청구 등록]  →  '유상(견적요청)' 모드 선택",
     "현장·호기를 먼저 선택하고, 교체할 자재(품목·수량)와 신청사유를 입력해 접수합니다.",
     "견적요청 접수(요청번호 QR-####) · 자재 담당자에게 알림 발송 · 상태 '신청'"),
    (2, AMBER, "견적 담당",
     "메뉴: [견적요청 목록]  →  해당 요청의 [견적서 작성]",
     "요청 내용이 자동으로 채워집니다. 자재 단가·인건비(공임)·특기사항을 입력해 견적서를 저장합니다.",
     "견적서 발급(견적번호 Q-####) · 요청과 자동 연결 · 요청 상태 '견적발행'"),
    (3, VIOLET, "관리자 (결재)",
     "메뉴: [견적서 목록] / [견적 상세]",
     "작성된 견적을 검토하고 결재합니다.  상태: 작성중 → 발행 → 승인.",
     "결재 '승인' 후에만 자재신청·세금계산서 발행이 가능해집니다."),
    (4, EMERALD, "자재 담당",
     "[견적 상세]에서 '자재신청 생성'  →  [자재신청 관리]에서 출고처리",
     "승인된 견적을 근거로 자재신청을 만들고, 재고 확인 후 출고 처리합니다.",
     "출고 기록 생성 · 재고 차감 · 진행상태 '자재신청' → '자재출고'"),
    (5, ROSE, "회계 / 견적 담당",
     "[견적 상세]  →  [🧾 세금계산서 발행] / [📄 거래명세서 발행]",
     "승인된 유상 견적에서 발행 버튼을 눌러 발행하고, 인쇄(PDF 저장)합니다.",
     "발행 이력 기록(번호 TX-#### / DN-####) · 발행시점 금액 보존 · 진행상태 '세금계산서발급'"),
    (6, NAVY, "회계",
     "[견적 상세]  →  [입금 등록]",
     "입금 내역을 등록합니다. 분할 입금·선입금도 지원합니다.",
     "누적 입금이 견적 합계 이상이면 진행상태 '입금완료' → 거래 '종료'"),
]
titles = ["견적요청 등록", "견적서 작성", "결재 승인", "자재 신청·출고", "세금계산서·거래명세서 발행", "입금 관리"]
for idx, (num, ac, who, where, what, result) in enumerate(details):
    s = new_slide()
    header_bar(s, f"STEP {num}. {titles[idx]}", "담당 · 화면 · 하는 일 · 시스템 결과", accent=ac)
    step_detail(s, num, ac, who, where, what, result)
    footer(s, 3 + idx, TOTAL)

# ─── 9. 진행상태 6단계 ───
s = new_slide()
header_bar(s, "견적 진행상태 (한 건의 일생)", "견적서마다 아래 6단계로 진행됩니다", accent=NAVY)
states = [
    ("미시작", "견적 작성/승인", SLATE_500),
    ("자재신청", "자재신청 생성", BLUE),
    ("자재출고", "출고 처리 완료", EMERALD),
    ("세금계산서발급", "세금계산서 발행", ROSE),
    ("입금완료", "대금 전액 입금", AMBER),
    ("종료", "거래 마감", SLATE_900),
]
n = len(states)
bw, bh = Inches(1.85), Inches(1.5)
gap = Inches(0.16)
total_w = bw * n + gap * (n - 1)
x0 = (WIDE_W - total_w) / 2
y0 = Inches(2.6)
for i, (title, desc, ac) in enumerate(states):
    x = x0 + i * (bw + gap)
    add_rect(s, x, y0, bw, bh, fill=WHITE, line=ac, line_w=1.5)
    add_text(s, x, y0 + Inches(0.25), bw, Inches(0.5), title, size=13, bold=True,
             color=ac if ac != SLATE_500 else SLATE_700, align=PP_ALIGN.CENTER)
    add_text(s, x, y0 + Inches(0.85), bw, Inches(0.5), desc, size=10, color=SLATE_500,
             align=PP_ALIGN.CENTER)
    if i < n - 1:
        arrow_right(s, x + bw - Inches(0.02), y0 + bh / 2 - Inches(0.13),
                    width=Inches(0.16), height=Inches(0.26), color=SLATE_500)
add_text(s, Inches(0.5), Inches(4.7), Inches(12.3), Inches(0.5),
         "진행상태는 각 작업(자재신청·출고·발행·입금) 시 자동으로 갱신됩니다.",
         size=12, color=SLATE_700, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(5.3), Inches(12.3), Inches(0.9),
         "견적 상세 화면 상단에서 현재 진행상태를 항상 확인할 수 있어, 어느 현장의 청구가 어디까지 진행됐는지 한눈에 파악됩니다.",
         size=12, color=SLATE_500, align=PP_ALIGN.CENTER)
footer(s, 9, TOTAL)

# ─── 10. 역할 / 권한 요약 ───
s = new_slide()
header_bar(s, "역할별 담당 정리", "누가 · 어느 단계를 · 어디서", accent=BLUE)
rows = [
    ("현장 보수원", "① 견적요청 등록", "견적 및 자재청구 등록", BLUE),
    ("견적 담당", "② 견적서 작성", "견적요청 목록 → 견적서", AMBER),
    ("관리자", "③ 결재 승인", "견적 상세", VIOLET),
    ("자재 담당", "④ 자재신청·출고", "견적 상세 / 자재신청 관리", EMERALD),
    ("회계 담당", "⑤ 세금계산서 발행 · ⑥ 입금", "견적 상세 → 발행/입금", ROSE),
]
ry = Inches(1.4)
rh = Inches(0.92)
for who, step, where, ac in rows:
    add_rect(s, Inches(0.7), ry, Inches(11.9), rh, fill=WHITE, line=SLATE_200, line_w=0.75)
    add_rect(s, Inches(0.7), ry, Inches(0.1), rh, fill=ac, line=None, shape=MSO_SHAPE.RECTANGLE)
    add_text(s, Inches(1.0), ry, Inches(2.6), rh, who, size=14, bold=True, color=SLATE_900, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(3.7), ry, Inches(4.6), rh, step, size=13, color=ac if ac != SLATE_500 else SLATE_700, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(8.4), ry, Inches(4.0), rh, where, size=12, color=SLATE_500, anchor=MSO_ANCHOR.MIDDLE)
    ry = ry + rh + Inches(0.12)
footer(s, 10, TOTAL)

# ─── 11. 꼭 기억하세요 ───
s = new_slide()
header_bar(s, "꼭 기억하세요", "자주 묻는 사항 · 주의점", accent=ROSE)
tips = [
    ("유상 vs 무상", "FM(무상) 현장 자재는 견적·세금계산서 대상이 아닙니다. '자재신청'으로 처리하세요.", AMBER),
    ("승인 후 발행", "세금계산서 발행·진행상태 변경은 결재 [승인] 이후에만 가능합니다.", VIOLET),
    ("발행 = 이력 + 인쇄", "시스템 발행은 발행 이력 기록과 사내 양식 인쇄(PDF)를 제공합니다. 국세청 실제 발행(홈택스/이카운트)은 별도 진행합니다.", BLUE),
    ("금액은 발행 시점 고정", "발행 후 견적을 수정해도 발행된 세금계산서 금액은 그대로 보존됩니다.", EMERALD),
    ("추적은 견적 단위", "견적서 한 건을 기준으로 신청·출고·발행·입금이 연결됩니다. 진행상태로 확인하세요.", NAVY),
]
ty = Inches(1.35)
th = Inches(1.0)
for title, body, ac in tips:
    add_rect(s, Inches(0.7), ty, Inches(11.9), th, fill=SLATE_50, line=SLATE_200, line_w=0.75)
    add_rect(s, Inches(0.7), ty, Inches(0.1), th, fill=ac, line=None, shape=MSO_SHAPE.RECTANGLE)
    add_text(s, Inches(1.0), ty + Inches(0.12), Inches(3.2), Inches(0.75), title, size=13, bold=True, color=ac if ac != SLATE_500 else SLATE_700, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(4.3), ty + Inches(0.1), Inches(8.1), th - Inches(0.2), body, size=12, color=SLATE_900, anchor=MSO_ANCHOR.MIDDLE)
    ty = ty + th + Inches(0.12)
footer(s, 11, TOTAL)

# ─── 12. 마무리 ───
s = new_slide()
add_rect(s, Inches(0), Inches(0), WIDE_W, WIDE_H, fill=NAVY, line=None, shape=MSO_SHAPE.RECTANGLE)
add_text(s, Inches(0.8), Inches(2.4), Inches(11.7), Inches(1.0),
         "한 줄 요약", size=20, color=RGBColor(0x93, 0xC5, 0xFD))
add_text(s, Inches(0.8), Inches(3.1), Inches(11.7), Inches(1.6),
         "견적요청 ▸ 견적서 ▸ 승인 ▸ 자재출고 ▸ 세금계산서 발행 ▸ 입금",
         size=26, bold=True, color=WHITE)
add_text(s, Inches(0.8), Inches(4.7), Inches(11.7), Inches(0.6),
         "모든 단계가 견적 한 건에 연결되어 진행상태로 추적됩니다.", size=15, color=SLATE_200)
add_text(s, Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.4),
         "DS_Material  ·  궁금한 점은 자재/회계 담당자에게 문의하세요.", size=12, color=SLATE_500)

# ── 저장 ──
out = Path(__file__).resolve().parent.parent / "docs" / "견적-세금계산서_발행안내_DS_Material.pptx"
out.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(out))
print(f"저장 완료: {out}  (슬라이드 {len(prs.slides._sldIdLst)}장)")
